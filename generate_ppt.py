#!/usr/bin/env python3
"""四地上市规则PPT生成"""

from pptx import Presentation
from pptx.util import Inches, Pt
# from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
import os

# 创建PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 标题样式
def set_title(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True

def add_bullet(slide, text, level=0):
    txBox = slide.shapes.placeholders[1]
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(18)

# 幻灯片1: 封面
slide_layout = prs.slide_layouts[6]  # 空白
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "北上深港四地上市规则全攻略"
p.font.size = Pt(44)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "投资人 · 被投企业 · 投行从业者必读"
p.font.size = Pt(24)
p.alignment = PP_ALIGN.CENTER

# 幻灯片2: 目录
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "目录")
content = slide.placeholders[1]
content.text = "一、上海证券交易所主板\n二、上海证券交易所科创板\n三、深圳证券交易所创业板\n四、香港联交所主板\n五、四地对比分析\n六、上市流程全图"

# 幻灯片3: 上交所主板
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "一、上海证券交易所主板")
content = slide.placeholders[1]
content.text = """法规依据：《上证发〔2025〕51号》《证监会令第205号》

发行条件：
• 持续经营满3年
• 会计基础规范，内控制度健全
• 业务完整，独立持续经营能力
• 经营合法合规

上市标准（三选一）：
• 标准一：净利润3年累计≥2亿，营收≥15亿
• 标准二：市值≥50亿，净利润为正，营收≥6亿
• 标准三：市值≥100亿，净利润为正，营收≥10亿"""

# 幻灯片4: 科创板
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "二、上海证券交易所科创板")
content = slide.placeholders[1]
content.text = """法规依据：《上证发〔2025〕60号》《科创属性评价指引》

科创板定位：
• 新一代信息技术、高端装备、新材料
• 新能源、生物医药、高端服务
• 人工智能、量子科技、脑机接口等

上市标准（五选一）：
• 标准一：市值≥10亿，营收≥1亿
• 标准二：市值≥15亿，研发累计≥5亿
• 标准三：市值≥20亿，研发累计≥2亿，核心技术
• 标准四：市值≥30亿，营收≥3亿
• 标准五：市值≥40亿，营收≥5亿"""

# 幻灯片5: 创业板
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "三、深圳证券交易所创业板")
content = slide.placeholders[1]
content.text = """法规依据：《深证上〔2025〕394号》

创业板定位：成长型创新创业企业

负面清单：农林牧渔、采矿业、酒饮料、纺织业、
         黑色金属、电力燃气、建筑业、批发零售等

上市标准（四选一）：
• 标准一：市值≥10亿，净利润2年累计≥1亿
• 标准二：市值≥10亿，净利润为正，营收≥4亿
• 标准三：市值≥30亿，营收≥6亿
• 标准四：市值≥50亿，营收≥8亿（未盈利）"""

# 幻灯片6: 港股主板
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "四、香港联交所主板")
content = slide.placeholders[1]
content.text = """法规依据：《主板上市规则》《公司条例》

上市标准（三选一）：
• 标准一：盈利≥5000万港币3年，市值≥5亿
• 标准二：市值≥40亿，营收≥6亿，现金流≥1亿
• 标准三：市值≥20亿，营收≥5亿，现金流≥1亿

特殊途径：
• 18A章：生物科技，核心产品通过一期临床
• 18C章：特专科技，市值≥80亿
• SPAC：特殊目的收购公司，最低集资10亿"""

# 幻灯片7: 四地对比
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "五、四地核心指标对比")
content = slide.placeholders[1]
content.text = """财务指标对比：
┌──────────┬────────┬───────┬──────┬────────┐
│ 市场     │ 最低市值│ 盈利要求│ 营收 │ 未盈利 │
├──────────┼────────┼───────┼──────┼────────┤
│ 上交所主板│ 50亿   │ 3年2亿 │ 15亿 │  不支持│
│ 科创板   │ 10亿   │ 1年正  │ 1亿  │  支持  │
│ 创业板   │ 10亿   │ 2年1亿 │ 无   │  支持  │
│ 港股主板 │ 5亿    │ 3年5千万│ 无   │  支持  │
└──────────┴────────┴───────┴──────┴────────┘

行业定位：
• 主板：成熟大型企业
• 科创板：科技创新
• 创业板：成长型创新
• 港股：国际化"""

# 幻灯片8: 上市流程
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "六、上市流程")
content = slide.placeholders[1]
content.text = """A股流程（注册制）：
1. 改制设立 → 2. 辅导验收(12月+) → 3. 申报审核
4. 注册发行 → 5. 上市交易
审核周期：3-12个月

港股流程：
1. 准备阶段 → 2. 审批(聆讯) → 3. 发行
4. 上市交易
审核周期：3-6个月

关键时间：
• A股：18-36个月
• 港股：12-24个月"""

# 幻灯片9: 建议
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "七、实务建议")
content = slide.placeholders[1]
content.text = """如何选择上市地？

• 传统行业、盈利稳定 → 上交所主板
• 科技创新、未盈利 → 科创板
• 成长型、创新业态 → 深交所创业板
• 国际化、美元基金 → 港交所主板
• 生物医药早期 → 港交所18A
• 特专科技 → 港交所18C

注意事项：
• A股：提前规划社保公积金合规
• 港股：保荐人资质重要，国际路演成本高
• 两地上市：可A+H同步"""

# 幻灯片10: 联系方式
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide, "附录：官方联系方式")
content = slide.placeholders[1]
content.text = """资料来源：

• 上海证券交易所：www.sse.com.cn
• 深圳证券交易所：www.szse.cn
• 香港联交所：www.hkex.com.hk
• 中国证监会：www.csrc.gov.cn

*本资料仅供参考，具体要求请以各交易所官方最新规则为准*

编制日期：2026年3月"""

# 保存
output = "/Users/henryt/.openclaw/workspace/reports/ipo-guide-四地上市规则.pptx"
prs.save(output)
print(f"PPT已保存: {output}")
